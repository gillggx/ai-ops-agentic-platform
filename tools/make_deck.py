# -*- coding: utf-8 -*-
"""
產出決策報告 PPT：
  Part 1 — LLM 選型：KIMI K2.5 vs GLM-5.2（驅動自研 builder agent）
  Part 2 — 建構模式：cowork（Claude Desktop + MCP）vs 自研 agent 優缺點
  Part 3 — 建議與決策矩陣
資料來源：docs/LLM_BAKEOFF_KIMI_vs_GLM.html + tools/slash17/BASELINE.md + 架構記憶
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
NAVY   = RGBColor(0x0B, 0x1F, 0x4D)
BLUE   = RGBColor(0x1D, 0x4E, 0xD8)
BLUEBG = RGBColor(0xEF, 0xF6, 0xFF)
INK    = RGBColor(0x1F, 0x29, 0x33)
INK2   = RGBColor(0x10, 0x18, 0x28)
MUT    = RGBColor(0x66, 0x70, 0x85)
LINE   = RGBColor(0xE4, 0xE7, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x06, 0x76, 0x47)
GREENBG= RGBColor(0xEC, 0xFD, 0xF3)
RED    = RGBColor(0xB4, 0x23, 0x18)
REDBG  = RGBColor(0xFE, 0xF3, 0xF2)
AMBER  = RGBColor(0xB5, 0x47, 0x08)
AMBERBG= RGBColor(0xFF, 0xFA, 0xEB)
GREYBG = RGBColor(0xF8, 0xFA, 0xFC)

FONT = "Microsoft JhengHei"   # 繁中

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def rect(s, l, t, w, h, fill, line=None, line_w=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def settext(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_sp=1.0):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", align)
        p.space_after = Pt(para.get("space_after", space_after))
        p.space_before = Pt(para.get("space_before", 0))
        p.line_spacing = para.get("line_sp", line_sp)
        for j, seg in enumerate(para["segs"]):
            r = p.add_run(); r.text = seg["t"]
            f = r.font
            f.name = FONT
            f.size = Pt(seg.get("sz", 14))
            f.bold = seg.get("b", False)
            f.italic = seg.get("i", False)
            f.color.rgb = seg.get("c", INK)


def seg(t, sz=14, b=False, c=INK, i=False):
    return {"t": t, "sz": sz, "b": b, "c": c, "i": i}


def para(segs, align=PP_ALIGN.LEFT, sa=4, sb=0, ls=1.0):
    return {"segs": segs, "align": align, "space_after": sa, "space_before": sb, "line_sp": ls}


def header_bar(s, kicker, title, sub=None):
    rect(s, 0, 0, 13.333, 1.18, NAVY)
    rect(s, 0, 1.18, 13.333, 0.06, BLUE)
    t = box(s, 0.55, 0.16, 12.2, 1.0)
    runs = [para([seg(kicker, 11, True, RGBColor(0x9D, 0xC0, 0xFF))], sa=2),
            para([seg(title, 23, True, WHITE)], sa=0)]
    if sub:
        runs.append(para([seg(sub, 12, False, RGBColor(0xCE, 0xDC, 0xF5))], sb=2))
    settext(t.text_frame, runs)


def footer(s, n):
    t = box(s, 0.55, 7.02, 12.2, 0.35)
    settext(t.text_frame, [para([
        seg("Pipeline Builder 智能化方案評估  ·  2026-06  ·  內部決策報告", 9, False, MUT),
    ])])
    p = box(s, 12.2, 7.02, 0.7, 0.35)
    settext(p.text_frame, [para([seg(str(n), 9, True, MUT)], align=PP_ALIGN.RIGHT)])


def chip(s, l, t, w, h, txt, fg, bg, bd, sz=11):
    sp = rect(s, l, t, w, h, bg, bd, 0.75)
    settext(sp.text_frame, [para([seg(txt, sz, True, fg)], align=PP_ALIGN.CENTER)],
            anchor=MSO_ANCHOR.MIDDLE)
    return sp


# simple table helper -------------------------------------------------------
def table(s, l, t, col_w, rows, header_fill=NAVY, header_fg=WHITE,
          row_h=0.34, fsz=11.5, hsz=11.5, zebra=True, aligns=None):
    """rows[0] = header. col_w: list inches. cell = str or (str, {opts})."""
    x = l
    total_w = sum(col_w)
    n = len(rows)
    cur_t = t
    for ri, row in enumerate(rows):
        x = l
        # row background
        if ri == 0:
            rect(s, l, cur_t, total_w, row_h, header_fill)
        elif zebra and ri % 2 == 0:
            rect(s, l, cur_t, total_w, row_h, GREYBG)
        for ci, cell in enumerate(row):
            opts = {}
            if isinstance(cell, tuple):
                txt, opts = cell
            else:
                txt = cell
            cw = col_w[ci]
            tb = box(s, x, cur_t, cw, row_h)
            if ri == 0:
                c = opts.get("c", header_fg); b = True; sz = hsz
            else:
                c = opts.get("c", INK); b = opts.get("b", False); sz = opts.get("sz", fsz)
            al = opts.get("align", PP_ALIGN.LEFT)
            settext(tb.text_frame, [para([seg(txt, sz, b, c)], align=al, ls=0.95)],
                    anchor=MSO_ANCHOR.MIDDLE)
            x += cw
        cur_t += row_h
    # bottom border lines
    for ri in range(n + 1):
        ln = rect(s, l, t + ri * row_h - 0.005, total_w, 0.012, LINE)
    return cur_t


# ============================================================ SLIDE 1 — TITLE
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.55, 13.333, 0.05, BLUE)
t = box(s, 0.9, 1.7, 11.5, 2.6)
settext(t.text_frame, [
    para([seg("LLM 選型 · 建構模式評估", 14, True, RGBColor(0x9D, 0xC0, 0xFF))], sa=10),
    para([seg("智能 Pipeline 建構方案", 40, True, WHITE)], sa=2),
    para([seg("KIMI vs GLM-5.2　＋　cowork vs 自研 Agent", 22, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0),
])
t2 = box(s, 0.9, 4.75, 11.5, 1.6)
settext(t2.text_frame, [
    para([seg("兩個決策、一份報告：", 13, True, WHITE)], sa=6),
    para([seg("Q1  哪個 LLM 驅動我們的 builder agent — KIMI K2.5 還是 GLM-5.2？", 13, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=3),
    para([seg("Q2  我們該繼續自研 agent，還是改用 cowork（Claude Desktop + MCP）建 pipeline？", 13, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0),
])
t3 = box(s, 0.9, 6.5, 11.5, 0.5)
settext(t3.text_frame, [para([seg("資料來源：SLASH-17 端到端驗證（2026-06-25）· Remote MCP e2e 驗證（2026-06-26）", 11, False, MUT)])])

# ====================================================== SLIDE 2 — EXEC SUMMARY
s = slide()
header_bar(s, "EXECUTIVE SUMMARY", "結論先講")
cards = [
    ("Q1 — LLM 選型", "GLM-5.2 勝出，已上 Prod", GREEN, GREENBG,
     ["寬鬆功能標準 17/17 vs KIMI 15/17", "跑完全套快 61%（28 vs 73 分）、少 29% 呼叫", "Prompt cache 40-58% vs KIMI 實測 0%", "切換點：sidecar .env 一行（OLLAMA_MODEL）"]),
    ("Q2 — 建構模式", "雙軌，不是二選一", BLUE, BLUEBG,
     ["自研 agent = 嵌入產品的『生產面』，flow 可控、可測、可治理", "cowork = 強推理的『共創 / authoring 面』，適合探索與打樣", "cowork 不可直接執行危險動作 → 已建 UI-handoff 橋接", "Fab 資料外送、不可測、vendor lock 是 cowork 的硬限制"]),
]
x = 0.55
for title, head, accent, bg, bullets in cards:
    w = 5.9
    rect(s, x, 1.55, w, 5.2, WHITE, LINE, 1.0)
    rect(s, x, 1.55, w, 0.12, accent)
    tb = box(s, x + 0.3, 1.8, w - 0.6, 0.9)
    settext(tb.text_frame, [
        para([seg(title, 12, True, MUT)], sa=3),
        para([seg(head, 19, True, accent)], sa=0),
    ])
    bb = box(s, x + 0.3, 2.9, w - 0.55, 3.7)
    runs = []
    for b in bullets:
        runs.append(para([seg("•  ", 14, True, accent), seg(b, 13.5, False, INK)], sa=12, ls=1.05))
    settext(bb.text_frame, runs)
    x += w + 0.43
footer(s, 2)

# ============================================ SLIDE 3 — PART 1 DIVIDER
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 1.0, 3.05, 1.4, 0.12, BLUE)
t = box(s, 1.0, 3.3, 11, 2)
settext(t.text_frame, [
    para([seg("PART 1", 16, True, RGBColor(0x9D, 0xC0, 0xFF))], sa=8),
    para([seg("LLM 選型：KIMI K2.5 vs GLM-5.2", 34, True, WHITE)], sa=6),
    para([seg("『哪個模型驅動我們自研的 builder agent』— 用 SLASH-17（17 條 production 指令）端到端驗證", 15, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0),
])
footer(s, 3)

# ==================================== SLIDE 4 — agent characteristics needs
s = slide()
header_bar(s, "PART 1 · 前提", "我們的 builder 是什麼樣的 agent",
           "graph-driven、逐 round 動作的 ReAct agent；對 LLM 的要求與一般 chat 不同 — 這決定了選型標準")
rows = [
    ["#", "對 LLM 的特性需求", "為什麼 critical"],
    ["1", "逐 round function calling", "每 round 只發 1 個 tool call；tool-call 不穩，整條建構崩"],
    ["2", "結構化 JSON 輸出", "goal_plan 出 3-7 phase JSON；parse 失敗就 retry / fail"],
    ["3", "長多步推理（≤32 round/phase）", "每 round 看 13 段大 observation（含整個 block catalog），要跨 round 連貫"],
    ["4", "中英混雜", "prompt + 機台指令雙語（中位數/median、偵測/找出 同義變形）"],
    ["5", "composite / fleet block 接線", "list_objects→mcp_foreach 全廠 fan-out、spc_panel 複合 — 最難推理"],
    ["6", "prompt-cache 友善（成本）", "大 observation prefix 每 round 重複；cache 命中率 = 成本主槓桿"],
    ["7", "果斷、不過度建構", "該用 minimal pipeline，別亂加 node（verbose catalog 會誘發 over-build）"],
    ["8", "自我修正", "回應 verifier reject、改參數 / 重接線，不要繞圈"],
]
table(s, 0.55, 1.55, [0.55, 4.0, 7.6], rows, row_h=0.585, fsz=12, hsz=11.5,
      aligns=None)
footer(s, 4)

# ==================================== SLIDE 5 — KIMI vs GLM fit table
s = slide()
header_bar(s, "PART 1 · 契合度", "KIMI vs GLM：對這些特性的契合")
def b_win(txt="GLM"): return (txt, {"align": PP_ALIGN.CENTER, "b": True, "c": BLUE})
def b_tie(txt="平"):  return (txt, {"align": PP_ALIGN.CENTER, "b": True, "c": MUT})
rows = [
    ["特性", "KIMI K2.5", "GLM-5.2 @ Fireworks", "勝"],
    ["1 · tool calling", "需 XML 解析補（native 較弱）", "native function-calling 紮實", b_win()],
    ["2 · JSON 輸出", "OK", "OK、更精簡（output token 少一半）", b_tie("平/GLM")],
    ["3 · 長多步推理", "較囉嗦（229 round/17）", "reasoning 模型，少 29%（163）", b_win()],
    ["4 · 中英混雜", "強（中文母語級）", "強（中文母語級）", b_tie()],
    ["5 · composite/fleet 接線", "會掙扎（系統弱點）", "會掙扎、且 run-to-run 抖", b_tie("平(都弱)")],
    ["6 · prompt-cache（成本）", "實測 cache = 0", "cache 40-58%", b_win()],
    ["7 · 果斷/不過度建構", "慢、call 多", "快、但偶 over-build", b_tie()],
    ["8 · 自我修正", "會做、慢", "會做、turnaround 快", b_win()],
    ["速度（builder UX）", "110-533s / 案", "28-142s / 案（快 2-4×）", b_win()],
]
end_t = table(s, 0.55, 1.5, [3.1, 3.7, 4.3, 1.1], rows, row_h=0.46, fsz=11.5, hsz=11)
footer(s, 5)

# ==================================== SLIDE 6 — SLASH17 metric cards
s = slide()
header_bar(s, "PART 1 · 結果", "SLASH-17 端到端驗證 — 四個關鍵數字",
           "17 條 production builder 指令、同一套設定、人工標註功能可用性（寬鬆 = 建成 + 型態正確 + 抓到核心意圖）")
metrics = [
    ("寬鬆通過（功能可用）", "17/17", "vs 15/17", "GLM 修掉 KIMI 兩個實質錯", GREEN),
    ("跑完全套耗時", "28 min", "vs 73 min", "GLM 快 61%（2.6×）", BLUE),
    ("總 LLM 呼叫", "163", "vs 229", "GLM 少 29%", BLUE),
    ("Prompt cache 命中", "40-58%", "vs 0%", "KIMI@Fireworks 實測 0", GREEN),
]
x = 0.55
w = 2.96
for lbl, val, vs, delta, accent in metrics:
    rect(s, x, 1.75, w, 2.5, WHITE, LINE, 1.0)
    rect(s, x, 1.75, w, 0.1, accent)
    tb = box(s, x + 0.22, 2.0, w - 0.4, 2.15)
    settext(tb.text_frame, [
        para([seg(lbl, 11, True, MUT)], sa=10, ls=1.0),
        para([seg(val, 30, True, INK2)], sa=2),
        para([seg(vs, 14, True, MUT)], sa=10),
        para([seg(delta, 12, True, accent)], sa=0, ls=1.05),
    ])
    x += w + 0.18
# key reading callout
cb = rect(s, 0.55, 4.5, 12.23, 1.95, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
settext(cb.text_frame, [
    para([seg("關鍵讀法", 13, True, INK2)], sa=6),
    para([seg("KIMI 兩個 failed 是型態 / 分析的實質錯：", 13, False, INK),
          seg("spc-xbar-r-pair", 13, True, RED), seg(" 畫了折線而非管制圖、", 13, False, INK),
          seg("apc-drift", 13, True, RED), seg(" 漏了 drift 判定 — GLM 這兩案都做對。", 13, False, INK)], sa=5, ls=1.1),
    para([seg("GLM 在嚴格 golden 標準下的失分（ooc-pareto / ooc-ranking）多為『缺排序』，功能上可用、且 KIMI 同樣中招。", 13, False, INK)], sa=0, ls=1.1),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# ==================================== SLIDE 7 — per-case detail (highlights)
s = slide()
header_bar(s, "PART 1 · 每案明細", "代表性案例：呼叫次數 / 耗時 / 寬鬆判定")
def ok():   return ("success", {"align": PP_ALIGN.CENTER, "b": True, "c": GREEN})
def fail(): return ("failed", {"align": PP_ALIGN.CENTER, "b": True, "c": RED})
def num(x): return (x, {"align": PP_ALIGN.RIGHT})
rows = [
    ["command", "KIMI 呼叫", "KIMI 耗時", "KIMI", "GLM 呼叫", "GLM 耗時", "GLM"],
    ["spc-xbar-r-pair", num("7"), num("119s"), fail(), num("5"), num("58s"), ok()],
    ["apc-drift", num("15"), num("403s"), fail(), num("14"), num("127s"), ok()],
    ["apc-trend（最極端）", num("40"), num("533s"), ok(), num("5"), num("44s"), ok()],
    ["apc-recipe-compare", num("20"), num("217s"), ok(), num("5"), num("55s"), ok()],
    ["spc-drift", num("13"), num("236s"), ok(), num("10"), num("51s"), ok()],
    ["ooc-ranking", num("18"), num("448s"), ok(), num("15"), num("120s"), ok()],
    ["patrol-status", num("29"), num("366s"), ok(), num("24"), num("205s"), ok()],
    ["spc-ooc（GLM 較慢案）", num("22"), num("263s"), ok(), num("21"), num("288s"), ok()],
    ["總計（17 案）", num("229"), num("73 min"), ("15/17", {"align": PP_ALIGN.CENTER, "b": True, "c": INK2}),
     num("163"), num("28 min"), ("17/17", {"align": PP_ALIGN.CENTER, "b": True, "c": BLUE})],
]
table(s, 0.55, 1.5, [3.0, 1.55, 1.55, 1.25, 1.55, 1.55, 1.25], rows, row_h=0.45, fsz=11.5, hsz=10.5)
note = box(s, 0.55, 6.55, 12.2, 0.5)
settext(note.text_frame, [para([seg("平均/案 — KIMI：13.5 call · 257s ｜ GLM：9.6 call · 100s。最極端 apc-trend：GLM 5 call/44s vs KIMI 40 call/533s。", 11, False, MUT)])])
footer(s, 7)

# ==================================== SLIDE 8 — Part 1 conclusion
s = slide()
header_bar(s, "PART 1 · 結論", "GLM-5.2 勝出 — 已切換 Prod")
items = [
    ("嚴格標準（golden 全中）", "兩者接近 ~14/17。KIMI 穩定，GLM 有 run-to-run 變異（12-15）。", AMBER, AMBERBG),
    ("寬鬆標準（功能可用）", "GLM 17/17 略勝 KIMI 15/17 — GLM 修掉 KIMI 兩個實質錯；GLM 失分多為 golden 過嚴或系統共通弱點（sort / fleet）。", GREEN, GREENBG),
    ("效率與成本", "GLM 少 29% 呼叫、跑完全套快 61%；cache 40-58% vs KIMI 實測 0% → 成本不輸還更穩。", BLUE, BLUEBG),
    ("待管理的唯一風險", "GLM 的 run-to-run 變異 → 以 hardening（phase 拆原子 / fleet hint / spc_panel fallback）壓低，這些對 KIMI 同樣有益。", AMBER, AMBERBG),
]
y = 1.6
for title, body, accent, bg in items:
    rect(s, 0.55, y, 12.23, 1.18, bg, accent, 1.0)
    tb = box(s, 0.85, y + 0.12, 11.7, 0.95)
    settext(tb.text_frame, [
        para([seg(title, 14, True, accent)], sa=3),
        para([seg(body, 13, False, INK)], sa=0, ls=1.05),
    ], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.32
dec = box(s, 0.55, y + 0.02, 12.23, 0.5)
settext(dec.text_frame, [para([
    seg("決策：", 13.5, True, INK2),
    seg("已切換 prod 至 GLM-5.2 @ Fireworks @ reasoning medium。切換點 = sidecar .env 一行（OLLAMA_MODEL + OPENROUTER_PROVIDER_ORDER），可一鍵回退。", 13, False, INK)])])
footer(s, 8)

# ============================================ SLIDE 9 — PART 2 DIVIDER
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 1.0, 3.05, 1.4, 0.12, BLUE)
t = box(s, 1.0, 3.3, 11.4, 2)
settext(t.text_frame, [
    para([seg("PART 2", 16, True, RGBColor(0x9D, 0xC0, 0xFF))], sa=8),
    para([seg("建構模式：cowork vs 自研 Agent", 32, True, WHITE)], sa=6),
    para([seg("『要不要改用 cowork（Claude Desktop + 我們的 Remote MCP）來建 pipeline，取代自研 builder？』", 15, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0),
])
footer(s, 9)

# ==================================== SLIDE 10 — two modes architecture
s = slide()
header_bar(s, "PART 2 · 兩種模式", "同一個目標、兩條路徑：誰來『推理建構』？")
# left card — self agent
lx = 0.55; w = 6.0
rect(s, lx, 1.55, w, 5.15, WHITE, LINE, 1.0)
rect(s, lx, 1.55, w, 0.5, BLUE)
hb = box(s, lx + 0.25, 1.6, w - 0.5, 0.42)
settext(hb.text_frame, [para([seg("A · 自研 Agent（現行 Prod）", 15, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
tb = box(s, lx + 0.3, 2.2, w - 0.6, 4.4)
settext(tb.text_frame, [
    para([seg("LLM：", 13, True, INK2), seg("GLM-5.2 @ Fireworks（我們選的）", 13, False, INK)], sa=7),
    para([seg("架構：", 13, True, INK2), seg("LangGraph 圖驅動 — goal_plan → agentic_phase_loop → phase_verifier", 13, False, INK)], sa=7, ls=1.05),
    para([seg("位置：", 13, True, INK2), seg("跑在我們 sidecar，內嵌 Glass Box GUI", 13, False, INK)], sa=7),
    para([seg("流程：", 13, True, INK2), seg("deterministic — 每步 graph node 決定，可單測、可 trace", 13, False, INK)], sa=7, ls=1.05),
    para([seg("資料：", 13, True, INK2), seg("留在我們基礎設施內，只送 prompt 給 LLM", 13, False, INK)], sa=0, ls=1.05),
])
# right card — cowork
rx = 6.78
rect(s, rx, 1.55, w, 5.15, WHITE, LINE, 1.0)
rect(s, rx, 1.55, w, 0.5, NAVY)
hb = box(s, rx + 0.25, 1.6, w - 0.5, 0.42)
settext(hb.text_frame, [para([seg("B · cowork（Claude Desktop + MCP）", 15, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
tb = box(s, rx + 0.3, 2.2, w - 0.6, 4.4)
settext(tb.text_frame, [
    para([seg("LLM：", 13, True, INK2), seg("Claude（使用者 Claude 訂閱裡的 Opus/Sonnet）", 13, False, INK)], sa=7, ls=1.05),
    para([seg("架構：", 13, True, INK2), seg("Claude 像人一樣用我們 MCP 6+11 個工具自由推理", 13, False, INK)], sa=7, ls=1.05),
    para([seg("位置：", 13, True, INK2), seg("Claude Desktop（外部 App），透過 Remote MCP 連回 EC2", 13, False, INK)], sa=7, ls=1.05),
    para([seg("流程：", 13, True, INK2), seg("Claude 自由意志 — 我們只能用 tool description 引導，不能強制 flow", 13, False, INK)], sa=7, ls=1.05),
    para([seg("資料：", 13, True, INK2), seg("block catalog + 資料 preview 送到 Anthropic 雲端", 13, False, INK)], sa=0, ls=1.05),
])
footer(s, 10)

# ==================================== SLIDE 11 — dimension comparison table
s = slide()
header_bar(s, "PART 2 · 維度對比", "九個維度的正面交鋒")
def win_a(): return ("自研 A", {"align": PP_ALIGN.CENTER, "b": True, "c": BLUE})
def win_b(): return ("cowork B", {"align": PP_ALIGN.CENTER, "b": True, "c": NAVY})
def tie():   return ("看情境", {"align": PP_ALIGN.CENTER, "b": True, "c": MUT})
rows = [
    ["維度", "自研 Agent (A)", "cowork (B)", "誰強"],
    ["模型推理力", "GLM-5.2 中上", "Claude 旗艦、明顯更強", win_b()],
    ["流程可控 / 可測", "graph 決定、可單測、有 BuildTracer", "Claude 自由意志、不可 trace 內部", win_a()],
    ["產品整合 / UX", "內嵌 Glass Box、即時 canvas", "外部 App、需切換、靠 handoff 橋接", win_a()],
    ["每次建構成本", "我們付 ~$0.1-0.4/案（OpenRouter）", "對我們 $0；但每人需 Claude 訂閱", tie()],
    ["治理 / 安全", "跑在我們 authz 內", "不可直接執行危險動作 → UI-handoff", win_a()],
    ["資料外送（Fab 關鍵）", "只送 prompt，可控", "互動 + 資料 preview 進 Anthropic 雲", win_a()],
    ["維運負擔", "整條 graph+prompt+選型都我們扛", "模型由 Anthropic 維護，我們只維 MCP 薄層", win_b()],
    ["難題（composite/fleet）", "會掙扎（系統弱點）", "Claude 較可能一次到位", win_b()],
    ["vendor lock / 可用性", "我們完全掌控", "綁 Anthropic 雲 + Desktop remote-MCP（仍在演進）", win_a()],
]
table(s, 0.55, 1.5, [2.7, 4.0, 4.3, 1.23], rows, row_h=0.50, fsz=11, hsz=11)
footer(s, 11)

# ==================================== SLIDE 12 — cowork pros / cons
s = slide()
header_bar(s, "PART 2 · cowork", "cowork（Claude Desktop + MCP）優缺點")
# pros
rect(s, 0.55, 1.55, 6.0, 5.15, GREENBG, RGBColor(0xA6, 0xF4, 0xC5), 1.0)
hb = box(s, 0.8, 1.7, 5.5, 0.4)
settext(hb.text_frame, [para([seg("優點", 16, True, GREEN)])])
pb = box(s, 0.85, 2.25, 5.5, 4.3)
pros = [
    "頂級推理力 — Claude 旗艦處理 composite / fleet 這類我們 agent 會掙扎的難題，較可能一次到位",
    "對我們零每案成本 — 算在使用者自己的 Claude 訂閱，不進我們 OpenRouter 帳單",
    "維運薄 — 模型由 Anthropic 維護升級，我們只維 MCP 工具層（6+11 工具）",
    "天生會探索 — 像人一樣 discover→preview→assemble→validate，適合打樣與共創",
    "升級免費搭便車 — Anthropic 換更強模型，我們不動 code 就受惠",
]
settext(pb.text_frame, [para([seg("＋  ", 14, True, GREEN), seg(p, 13, False, INK)], sa=11, ls=1.08) for p in pros])
# cons
rect(s, 6.78, 1.55, 6.0, 5.15, REDBG, RGBColor(0xFE, 0xCA, 0xCA), 1.0)
hb = box(s, 7.03, 1.7, 5.5, 0.4)
settext(hb.text_frame, [para([seg("缺點 / 風險", 16, True, RED)])])
cb = box(s, 7.08, 2.25, 5.5, 4.3)
cons = [
    "資料外送 — block catalog + 資料 preview 進 Anthropic 雲；半導體 Fab 的 data egress / air-gap 是硬限制",
    "流程不可控 — Claude 自由意志，我們只能用 tool desc 引導，無法強制 flow（違反『flow 寫 graph 不寫 prompt』原則）",
    "不可測 / 不可 trace — 沒有 BuildTracer、無法做 17-case 回歸；出錯只能猜",
    "需危險動作橋接 — 刪/停/啟用 rule 不能由 MCP 執行 → 必須走我們的 UI-handoff 人工確認",
    "vendor lock + 成熟度 — 綁 Anthropic 雲、Desktop remote-MCP 認證仍只到 secret-path；每位操作員要 Claude 帳號",
]
settext(cb.text_frame, [para([seg("－  ", 14, True, RED), seg(c, 13, False, INK)], sa=11, ls=1.08) for c in cons])
footer(s, 12)

# ==================================== SLIDE 13 — self agent pros / cons
s = slide()
header_bar(s, "PART 2 · 自研 Agent", "自研 Agent（GLM-5.2 builder）優缺點")
rect(s, 0.55, 1.55, 6.0, 5.15, GREENBG, RGBColor(0xA6, 0xF4, 0xC5), 1.0)
hb = box(s, 0.8, 1.7, 5.5, 0.4)
settext(hb.text_frame, [para([seg("優點", 16, True, GREEN)])])
pb = box(s, 0.85, 2.25, 5.5, 4.3)
pros = [
    "流程完全可控 — graph node 決定每一步，pure function 可單測；出錯知道是哪個 node",
    "可測 / 可回歸 — BuildTracer + trace_replay + verify-build；17-case 端到端基準隨時重跑",
    "深度產品整合 — 內嵌 Glass Box GUI、即時 canvas，使用者不離開系統",
    "資料留內網 — 只送 prompt 給 LLM，符合 Fab 資料治理；模型可換、甚至未來自架",
    "治理內建 — 跑在我們 authz / role 內，危險動作天然受控",
]
settext(pb.text_frame, [para([seg("＋  ", 14, True, GREEN), seg(p, 13, False, INK)], sa=12, ls=1.08) for p in pros])
rect(s, 6.78, 1.55, 6.0, 5.15, REDBG, RGBColor(0xFE, 0xCA, 0xCA), 1.0)
hb = box(s, 7.03, 1.7, 5.5, 0.4)
settext(hb.text_frame, [para([seg("缺點 / 限制", 16, True, RED)])])
cb = box(s, 7.08, 2.25, 5.5, 4.3)
cons = [
    "推理力天花板 — GLM-5.2 在 composite / fleet 難題會掙扎，且有 run-to-run 變異",
    "維運全包 — graph + prompt + 選型 + cost tuning 都要我們自己扛與調",
    "每案有成本 — ~$0.1-0.4/案進我們 OpenRouter 帳單（雖遠低於 cowork 的 per-seat）",
    "模型升級要自己追 — 換更強模型需重跑 bake-off、重調 prompt / hardening",
    "難題需持續 hardening — fleet hint、spc_panel fallback、phase 拆原子等要人力投入",
]
settext(cb.text_frame, [para([seg("－  ", 14, True, RED), seg(c, 13, False, INK)], sa=12, ls=1.08) for c in cons])
footer(s, 13)

# ==================================== SLIDE 14 — cost model
s = slide()
header_bar(s, "PART 2 · 成本", "成本模型不同維度 — per-build vs per-seat")
rows = [
    ["", "自研 Agent (A)", "cowork (B)"],
    ["計費對象", "我們（OpenRouter / Fireworks）", "使用者自己的 Claude 訂閱"],
    ["每次建構", "~$0.1 – $0.4 / 案", "對我們 $0（算在對方訂閱額度）"],
    ["固定成本", "0 額外席位；模型隨用隨付", "每位操作員需 Claude Pro/Team 席位（~$20-30/人/月）"],
    ["規模化拐點", "用量越大、總額線性成長但單價低", "人數越多、席位成本越高；少數 power-user 才划算"],
    ["成本可預測性", "高 — 我們掌控模型與 cache（命中 40-58%）", "低 — 取決於對方 plan 與 Anthropic 定價"],
]
table(s, 0.55, 1.6, [2.8, 5.0, 4.43], rows, row_h=0.62, fsz=12, hsz=12)
cb = rect(s, 0.55, 5.55, 12.23, 1.05, AMBERBG, RGBColor(0xFD, 0xE6, 0x8A), 1.0)
settext(cb.text_frame, [
    para([seg("讀法：", 13, True, AMBER), seg("兩者不是同一種錢。自研是『我們付、單價低、隨量成長』；cowork 是『使用者付、我們零邊際成本，但要人手一個 Claude 席位』。", 13, False, INK)], ls=1.1),
    para([seg("→ 全廠操作員規模下 per-seat 會超過 per-build；少數工程師打樣 / 共創時 cowork 反而便宜。", 13, True, INK2)], sb=3),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 14)

# ==================================== SLIDE 15 — governance / handoff
s = slide()
header_bar(s, "PART 2 · 治理", "為什麼 cowork 不能直接動產線 — 我們的 UI-handoff 橋接",
           "cowork 可以『建議』，但刪 / 停 / 啟用 rule 這類危險動作，必須回到我們系統由人確認")
steps = [
    ("1", "cowork 提議", "Claude 透過 MCP 呼叫 rule_request_delete / disable / activate", NAVY),
    ("2", "回傳 PENDING", "工具回 executed=false、PENDING_USER_CONFIRMATION + launch_url；不執行任何事", AMBER),
    ("3", "開我們 GUI", "使用者點開獨立確認頁（fit-window、無 app shell），看到影響範圍", BLUE),
    ("4", "人工確認執行", "在我們 authz 內、以登入身份執行 → 才真的刪/停/啟用", GREEN),
]
x = 0.55; w = 2.96
for n, title, body, accent in steps:
    rect(s, x, 1.85, w, 3.4, WHITE, LINE, 1.0)
    rect(s, x, 1.85, w, 0.7, accent)
    nb = box(s, x, 1.9, w, 0.6)
    settext(nb.text_frame, [para([seg(n, 24, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    tb = box(s, x + 0.22, 2.75, w - 0.44, 2.3)
    settext(tb.text_frame, [
        para([seg(title, 14, True, accent)], sa=6, align=PP_ALIGN.CENTER),
        para([seg(body, 12, False, INK)], sa=0, ls=1.1),
    ])
    if n != "4":
        ar = box(s, x + w - 0.05, 3.2, 0.35, 0.6)
        settext(ar.text_frame, [para([seg("→", 22, True, MUT)], align=PP_ALIGN.CENTER)])
    x += w + 0.18
cb = rect(s, 0.55, 5.5, 12.23, 1.1, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
settext(cb.text_frame, [
    para([seg("設計原則：", 13, True, INK2), seg("外部 Claude 永遠只是『提案者』，產線級動作的『執行權』鎖在我們 authenticated GUI 內（V63 ui_handoffs）。", 13, False, INK)], ls=1.1),
    para([seg("這同時補上了 cowork『流程不可控』的缺口 — 危險路徑被 deterministic 地擋在系統邊界。", 13, False, INK)], sb=3, ls=1.1),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 15)

# ============================================ SLIDE 16 — PART 3 DIVIDER
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 1.0, 3.05, 1.4, 0.12, BLUE)
t = box(s, 1.0, 3.3, 11, 2)
settext(t.text_frame, [
    para([seg("PART 3", 16, True, RGBColor(0x9D, 0xC0, 0xFF))], sa=8),
    para([seg("建議與決策矩陣", 34, True, WHITE)], sa=6),
    para([seg("不是二選一 — 兩條路各自最適的戰場", 15, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0),
])
footer(s, 16)

# ==================================== SLIDE 17 — recommendation (dual track)
s = slide()
header_bar(s, "PART 3 · 建議", "雙軌策略 — 各司其職")
# A
rect(s, 0.55, 1.6, 6.0, 4.0, WHITE, LINE, 1.0)
rect(s, 0.55, 1.6, 6.0, 0.55, BLUE)
hb = box(s, 0.8, 1.66, 5.5, 0.45)
settext(hb.text_frame, [para([seg("自研 Agent → 生產面（預設）", 14.5, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
tb = box(s, 0.85, 2.35, 5.45, 3.1)
settext(tb.text_frame, [
    para([seg("給：", 13, True, INK2), seg("全廠操作員、ON-DUTY、日常建 pipeline / rule", 13, False, INK)], sa=8, ls=1.05),
    para([seg("為何：", 13, True, INK2), seg("可控、可測、可治理、資料留內網、規模化單價低", 13, False, INK)], sa=8, ls=1.05),
    para([seg("模型：", 13, True, INK2), seg("GLM-5.2 @ Fireworks（Part 1 結論）", 13, False, INK)], sa=0, ls=1.05),
])
# B
rect(s, 6.78, 1.6, 6.0, 4.0, WHITE, LINE, 1.0)
rect(s, 6.78, 1.6, 6.0, 0.55, NAVY)
hb = box(s, 7.03, 1.66, 5.5, 0.45)
settext(hb.text_frame, [para([seg("cowork → 共創 / 打樣面", 14.5, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
tb = box(s, 7.08, 2.35, 5.45, 3.1)
settext(tb.text_frame, [
    para([seg("給：", 13, True, INK2), seg("少數工程師 / PE 探索新 pipeline、難題打樣", 13, False, INK)], sa=8, ls=1.05),
    para([seg("為何：", 13, True, INK2), seg("頂級推理一次到位 composite/fleet、對我們零邊際成本", 13, False, INK)], sa=8, ls=1.05),
    para([seg("護欄：", 13, True, INK2), seg("危險動作一律走 UI-handoff，產線執行權鎖在系統內", 13, False, INK)], sa=0, ls=1.05),
])
cb = rect(s, 0.55, 5.85, 12.23, 0.85, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
settext(cb.text_frame, [para([
    seg("一句話：", 13.5, True, INK2),
    seg("自研 agent 是『內嵌、可控、可測』的生產引擎；cowork 是『強推理、零邊際成本』的共創沙盒。兩者共用同一套 block / pipeline / handoff 基礎設施，互補不互斥。", 13, False, INK)], ls=1.15)],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 17)

# ==================================== SLIDE 18 — decision matrix
s = slide()
header_bar(s, "PART 3 · 決策矩陣", "什麼情境用哪個")
def A(): return ("自研 Agent", {"align": PP_ALIGN.CENTER, "b": True, "c": BLUE})
def B(): return ("cowork", {"align": PP_ALIGN.CENTER, "b": True, "c": NAVY})
def BOTH(): return ("兩者皆可", {"align": PP_ALIGN.CENTER, "b": True, "c": MUT})
rows = [
    ["情境", "選擇", "理由"],
    ["全廠操作員日常建 pipeline / rule", A(), "規模化、可控、可治理、資料留內網"],
    ["產線級危險動作（刪/停/啟用 rule）", A(), "執行權必須鎖在 authenticated GUI（即使 cowork 提案）"],
    ["需要回歸測試 / SLA 保證的流程", A(), "只有自研有 BuildTracer + 17-case 基準"],
    ["探索全新 / 複雜 composite·fleet pipeline", B(), "Claude 旗艦推理較可能一次到位"],
    ["工程師快速打樣、共創、catalog 探索", B(), "互動式探索是 cowork 的天生強項，零邊際成本"],
    ["嚴格 air-gap / 禁資料外送環境", A(), "cowork 會把 catalog + preview 送上 Anthropic 雲"],
    ["跨團隊知識沈澱（block 文件即文件）", BOTH(), "兩者都讀同一份 DB block description"],
]
table(s, 0.55, 1.55, [5.2, 1.9, 5.13], rows, row_h=0.62, fsz=12, hsz=12)
footer(s, 18)

# ==================================== SLIDE 19 — risks / next
s = slide()
header_bar(s, "PART 3 · 後續", "風險與下一步")
cols = [
    ("待管理風險", AMBER, AMBERBG, [
        "GLM run-to-run 變異 → hardening（phase 拆原子 / fleet hint / spc_panel fallback）",
        "cowork 資料外送 → 對 Fab 客戶須明確界定可送 / 不可送的資料範圍",
        "Remote MCP 認證僅 secret-path → 待 Claude Desktop 認證模型成熟升級 OAuth/Bearer",
        "cowork『流程不可控』→ 靠 UI-handoff 把危險路徑擋在邊界，非 prompt 約束",
    ]),
    ("下一步", BLUE, BLUEBG, [
        "Part 1：維持 GLM-5.2 @ Fireworks，持續以 SLASH-17 監測變異",
        "Part 2：cowork 定位為內部 power-user 工具，不開放全廠 / 不接觸產線執行",
        "成本：每季用 bake-off harness 重評選型（tools/slash17）",
        "治理：完成 cowork 資料外送白名單 + handoff 稽核軌跡",
    ]),
]
x = 0.55; w = 6.0
for title, accent, bg, items in cols:
    rect(s, x, 1.6, w, 5.0, bg, accent, 1.0)
    hb = box(s, x + 0.25, 1.75, w - 0.5, 0.45)
    settext(hb.text_frame, [para([seg(title, 16, True, accent)])])
    ib = box(s, x + 0.3, 2.35, w - 0.6, 4.0)
    settext(ib.text_frame, [para([seg("•  ", 14, True, accent), seg(it, 13, False, INK)], sa=14, ls=1.1) for it in items])
    x += w + 0.43
footer(s, 19)

out = "docs/PIPELINE_BUILDER_LLM_vs_COWORK.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
