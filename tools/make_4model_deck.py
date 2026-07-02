# -*- coding: utf-8 -*-
"""4-model builder LLM bake-off deck — one slide per model, each listing all 17 cases.
Data: KIMI/GLM (2026-06-25 bake-off) + Sonnet5/Opus4.8 (2026-07-01, trace-measured)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x0B,0x1F,0x4D); BLUE=RGBColor(0x1D,0x4E,0xD8); INK=RGBColor(0x1F,0x29,0x33)
INK2=RGBColor(0x10,0x18,0x28); MUT=RGBColor(0x66,0x70,0x85); LINE=RGBColor(0xE4,0xE7,0xEC)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0x06,0x76,0x47); GREENBG=RGBColor(0xEC,0xFD,0xF3)
RED=RGBColor(0xB4,0x23,0x18); REDBG=RGBColor(0xFE,0xF3,0xF2); AMBER=RGBColor(0xB5,0x47,0x08)
AMBERBG=RGBColor(0xFF,0xFA,0xEB); PURP=RGBColor(0x6D,0x28,0xD9); PURPBG=RGBColor(0xF5,0xF3,0xFF)
GREYBG=RGBColor(0xF8,0xFA,0xFC); SLATE=RGBColor(0x94,0xA3,0xB8); SKY=RGBColor(0x9D,0xC0,0xFF)
FONT="Microsoft JhengHei"; MONO="Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def box(s,l,t,w,h): return s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
def rect(s,l,t,w,h,fill,line=None,lw=0.75):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def seg(t,sz=14,b=False,c=INK,mono=False): return {"t":t,"sz":sz,"b":b,"c":c,"mono":mono}
def para(segs,align=PP_ALIGN.LEFT,sa=3,ls=1.0): return {"segs":segs,"align":align,"sa":sa,"ls":ls}
def settext(tf,runs,anchor=MSO_ANCHOR.TOP):
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,p in enumerate(runs):
        pp=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        pp.alignment=p.get("align",PP_ALIGN.LEFT); pp.space_after=Pt(p.get("sa",3)); pp.line_spacing=p.get("ls",1.0)
        for sg in p["segs"]:
            r=pp.add_run(); r.text=sg["t"]; f=r.font
            f.name=MONO if sg.get("mono") else FONT; f.size=Pt(sg.get("sz",14))
            f.bold=sg.get("b",False); f.color.rgb=sg.get("c",INK)
def chip(s,l,t,w,h,txt,fg,bg,bd,sz=10.5):
    sp=rect(s,l,t,w,h,bg,bd,0.75)
    settext(sp.text_frame,[para([seg(txt,sz,True,fg)],align=PP_ALIGN.CENTER)],anchor=MSO_ANCHOR.MIDDLE)

# case order
CASES=["spc-trend","spc-ooc","spc-cpk","spc-multi-tool","spc-drift","spc-xbar-r-pair",
"spc-multi-step","spc-tool-box","spc-normality","spc-cusum","apc-drift","apc-trend",
"apc-recipe-compare","patrol-status","ooc-ranking","ooc-pareto","step-yield"]
# per-model: metric label (calls or rounds), value, time_s, pass(bool)
KIMI={"spc-trend":(6,247,1),"spc-ooc":(22,263,1),"spc-cpk":(8,200,1),"spc-multi-tool":(8,216,1),
"spc-drift":(13,236,1),"spc-xbar-r-pair":(7,119,0),"spc-multi-step":(11,324,1),"spc-tool-box":(6,138,1),
"spc-normality":(6,125,1),"spc-cusum":(6,226,1),"apc-drift":(15,403,0),"apc-trend":(40,533,1),
"apc-recipe-compare":(20,217,1),"patrol-status":(29,366,1),"ooc-ranking":(18,448,1),"ooc-pareto":(8,191,1),"step-yield":(6,110,1)}
GLM={"spc-trend":(5,47,1),"spc-ooc":(21,288,1),"spc-cpk":(6,218,1),"spc-multi-tool":(5,86,1),
"spc-drift":(10,51,1),"spc-xbar-r-pair":(5,58,1),"spc-multi-step":(16,142,1),"spc-tool-box":(10,28,1),
"spc-normality":(6,40,1),"spc-cusum":(6,51,1),"apc-drift":(14,127,1),"apc-trend":(5,44,1),
"apc-recipe-compare":(5,55,1),"patrol-status":(24,205,1),"ooc-ranking":(15,120,1),"ooc-pareto":(5,95,1),"step-yield":(5,42,1)}
SON={"spc-trend":(5,33,1),"spc-ooc":(21,116,1),"spc-cpk":(36,306,1),"spc-multi-tool":(10,62,1),
"spc-drift":(12,66,1),"spc-xbar-r-pair":(9,57,0),"spc-multi-step":(14,115,1),"spc-tool-box":(6,40,1),
"spc-normality":(15,102,1),"spc-cusum":(7,46,1),"apc-drift":(11,104,0),"apc-trend":(3,26,1),
"apc-recipe-compare":(36,24,1),"patrol-status":(17,134,1),"ooc-ranking":(7,69,1),"ooc-pareto":(14,126,1),"step-yield":(6,39,1)}
OPUS={"spc-trend":(12,135,1),"spc-ooc":(24,169,1),"spc-cpk":(10,101,1),"spc-multi-tool":(11,93,1),
"spc-drift":(12,115,1),"spc-xbar-r-pair":(8,79,1),"spc-multi-step":(6,50,1),"spc-tool-box":(7,61,1),
"spc-normality":(7,52,1),"spc-cusum":(8,77,1),"apc-drift":(8,78,0),"apc-trend":(3,35,1),
"apc-recipe-compare":(12,39,1),"patrol-status":(16,127,1),"ooc-ranking":(6,71,1),"ooc-pareto":(7,71,1),"step-yield":(7,58,1)}
# 第 5 種:同一顆 GLM-5.2、跑在 Multi-Agent Phase 0 重構後的 builder(2026-07-02 p0_gate)
P0={"spc-trend":(6,25,1),"spc-ooc":(30,116,1),"spc-cpk":(34,284,1),"spc-multi-tool":(12,50,1),
"spc-drift":(17,79,1),"spc-xbar-r-pair":(9,58,1),"spc-multi-step":(8,31,1),"spc-tool-box":(7,113,1),
"spc-normality":(9,29,1),"spc-cusum":(9,46,1),"apc-drift":(12,134,1),"apc-trend":(3,17,1),
"apc-recipe-compare":(34,19,1),"patrol-status":(39,158,1),"ooc-ranking":(6,73,1),"ooc-pareto":(7,56,1),"step-yield":(9,49,1)}

def header(s,kicker,title,accent,sub=None):
    rect(s,0,0,13.333,1.12,NAVY); rect(s,0,1.12,13.333,0.06,accent)
    t=box(s,0.55,0.12,12.2,0.95)
    runs=[para([seg(kicker,11,True,SKY)],sa=2),para([seg(title,22,True,WHITE)],sa=0)]
    if sub: runs.append(para([seg(sub,11.5,False,RGBColor(0xCE,0xDC,0xF5))],sa=0))
    settext(t.text_frame,runs)
def footer(s,n):
    settext(box(s,0.55,7.04,11,0.33).text_frame,[para([seg("Builder LLM 四方 bake-off · SLASH-17 · 2026-07",9,False,MUT)])])
    settext(box(s,12.2,7.04,0.7,0.33).text_frame,[para([seg(str(n),9,True,MUT)],align=PP_ALIGN.RIGHT)])

def model_slide(n,name,data,accent,accentbg,metric_word,metrics,verdict,vfg,vbg,vbd,take):
    s=slide(); header(s,"每案結果 · "+metric_word.upper(),name,accent)
    # left metric panel
    rect(s,0.55,1.45,3.65,5.15,accentbg,accent,1.0)
    tb=box(s,0.78,1.68,3.2,4.7)
    runs=[para([seg("關鍵指標",13,True,accent)],sa=8)]
    for k,v in metrics:
        runs.append(para([seg(k+"  ",12,False,MUT),seg(v,15,True,INK2)],sa=8))
    settext(tb.text_frame,runs)
    chip(s,0.78,5.55,3.2,0.5,verdict,vfg,vbg,vbd,12)
    settext(box(s,0.78,6.08,3.2,0.5).text_frame,[para([seg(take,10.5,False,INK)],ls=1.05)])
    # right: 17 cases in 2 columns
    col_x=[4.45,8.62]; col_w=4.05
    for ci in range(2):
        rows=CASES[ci*9:ci*9+9] if ci==0 else CASES[9:]
        y=1.5
        for idx,case in enumerate(rows):
            gi=ci*9+idx
            val,tsec,ok=data[case]
            rh=0.555
            bg=WHITE if idx%2 else GREYBG
            rect(s,col_x[ci],y,col_w,rh,bg,LINE,0.5)
            # index+case
            settext(box(s,col_x[ci]+0.08,y,2.2,rh).text_frame,
                [para([seg(f"{gi+1:>2} ",9,True,MUT,mono=True),seg(case,10.5,False,INK)])],anchor=MSO_ANCHOR.MIDDLE)
            # metric+time
            settext(box(s,col_x[ci]+2.05,y,1.1,rh).text_frame,
                [para([seg(f"{val}{metric_word[0]} · {tsec}s",9.5,False,MUT,mono=True)],align=PP_ALIGN.RIGHT)],anchor=MSO_ANCHOR.MIDDLE)
            # status
            if ok: chip(s,col_x[ci]+col_w-0.66,y+0.12,0.55,0.31,"ok",GREEN,GREENBG,RGBColor(0xA6,0xF4,0xC5),9)
            else:  chip(s,col_x[ci]+col_w-0.66,y+0.12,0.55,0.31,"✗",RED,REDBG,RGBColor(0xFE,0xCA,0xCA),9)
            y+=rh
    footer(s,n)

# ---- 1 title
s=slide(); rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.6,13.333,0.05,BLUE)
settext(box(s,0.9,1.7,11.6,2.6).text_frame,[
    para([seg("Builder Agent · LLM 選型 · 四方彙總",14,True,SKY)],sa=10),
    para([seg("KIMI · GLM-5.2 · Sonnet 5 · Opus 4.8 · MA-P0",32,True,WHITE)],sa=4),
    para([seg("SLASH-17 端到端 bake-off — 每個模型各一頁完整結果",18,False,RGBColor(0xCE,0xDC,0xF5))],sa=0)])
settext(box(s,0.9,4.8,11.6,1.4).text_frame,[
    para([seg("結論:",13,True,WHITE),seg("GLM-5.2 = cost-right 預設(品質接近頂級、成本最低、cache 唯一有效)",13,False,RGBColor(0xCE,0xDC,0xF5))],sa=4),
    para([seg("Opus 4.8 品質天花板但 ~20× 成本(direct+cache) · Sonnet 5 被 GLM 支配 · KIMI 已淘汰",13,False,RGBColor(0xCE,0xDC,0xF5))],sa=0)])

# ---- 2 summary
s=slide(); header(s,"總表","四方一眼對照",BLUE)
hdr=["指標","KIMI K2.5","GLM-5.2","Sonnet 5","Opus 4.8","GLM@MA-P0"]
rows=[
 ("嚴格通過","~14/17","~14/17","15/17","16/17","17/17"),
 ("寬鬆(功能可用)","15/17","17/17","~17/17","~17/17","17/17"),
 ("跑完全套","73 min","28 min","24.4 min","23.5 min","22.3 min"),
 ("LLM 呼叫/round","229","163","229","164","237"),
 ("Prompt cache","0%","40–58%","63%*","52%*","54.7%"),
 ("公平成本/17*","~$3–4","~$1–3","$11.6*","$52*","~$1.3"),
 ("相對 GLM","~1.5×","1×","~5×","~20×","1×"),
 ("定位","淘汰","預設","無角色","升級層","現行 prod 架構"),
]
cw=[2.2,2.0,2.0,2.0,2.0,2.0]; x0=0.55; y0=1.55; rh=0.585
tot=sum(cw)
# header row
rect(s,x0,y0,tot,rh,NAVY); xx=x0
accents=[INK2,SLATE,GREEN,BLUE,PURP]
for ci,h in enumerate(hdr):
    settext(box(s,xx+0.05,y0,cw[ci]-0.1,rh).text_frame,[para([seg(h,12,True,WHITE)],align=PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)],anchor=MSO_ANCHOR.MIDDLE); xx+=cw[ci]
y=y0+rh
for ri,row in enumerate(rows):
    rect(s,x0,y,tot,rh,GREYBG if ri%2 else WHITE,LINE,0.5); xx=x0
    for ci,val in enumerate(row):
        c=INK; b=False
        if ci==0: b=True
        settext(box(s,xx+0.05,y,cw[ci]-0.1,rh).text_frame,[para([seg(val,11.5,b,c)],align=PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)],anchor=MSO_ANCHOR.MIDDLE); xx+=cw[ci]
    y+=rh
settext(box(s,0.55,y+0.1,12.2,1.1).text_frame,[
    para([seg("真正分高下只有 2 案:",12.5,True,INK2),seg("xbar_r(只有 GLM+Opus 建對管制圖)、weco_rules(四家幾乎都漏 = 系統層弱點)。其餘 15 案四模型都對。",12.5,False,INK)],ls=1.15),
    para([seg("* Sonnet/Opus 成本與 cache 為 2026-07-02 direct Anthropic + cache 補測(公平值);bakeoff 當日走 OpenRouter 無 cache 為 $18.7 / $68.6。",10.5,False,MUT)],sa=4,ls=1.1)])
footer(s,2)

# ---- 3-6 model slides
model_slide(3,"KIMI K2.5",KIMI,SLATE,RGBColor(0xF1,0xF5,0xF9),"calls",
    [("嚴格","~14/17"),("寬鬆","15/17"),("跑完","73 min"),("呼叫","229"),("cache","0%"),("成本","~$3–4")],
    "已淘汰",RED,REDBG,RGBColor(0xFE,0xCA,0xCA),
    "最慢(73min)、cache 0%;2 案失敗(xbar_r/apc-drift)。2026-06-25 已被 GLM 取代。")
model_slide(4,"GLM-5.2  （prod 預設）",GLM,GREEN,GREENBG,"calls",
    [("嚴格","~14/17"),("寬鬆","17/17"),("跑完","28 min"),("呼叫","163"),("cache","40–58%"),("成本","~$1–3")],
    "prod 預設 [best]",WHITE,GREEN,GREEN,
    "寬鬆 17/17、xbar_r 建對、cache 唯一有效、成本最低。CP 值完勝,現行 prod。")
model_slide(5,"Sonnet 5",SON,BLUE,RGBColor(0xEF,0xF6,0xFF),"rounds",
    [("嚴格","15/17"),("寬鬆","~17/17"),("跑完","24.4 min"),("round","229"),("cache","63%*"),("成本","$11.6*")],
    "被 GLM 支配",AMBER,AMBERBG,RGBColor(0xFD,0xE6,0x8A),
    "嚴格 15/17(xbar_r 退化成折線)、direct+cache 成本 ~5× GLM。更貴又不更準,無採用理由。")
model_slide(6,"Opus 4.8",OPUS,PURP,PURPBG,"rounds",
    [("嚴格","16/17"),("寬鬆","~17/17"),("跑完","23.5 min"),("round","164"),("cache","52%*"),("成本","$52*")],
    "品質天花板 · 升級層",PURP,PURPBG,RGBColor(0xC4,0xB5,0xFD),
    "嚴格 16/17 最高、最果斷(164 round),只漏 weco_rules。direct+cache 實測 ~20× GLM 成本。")
model_slide(7,"GLM-5.2 @ Multi-Agent P0  （第 5 種：新架構）",P0,GREEN,GREENBG,"rounds",
    [("嚴格","17/17"),("寬鬆","17/17"),("跑完","22.3 min"),("round","237"),("cache","54.7%"),("成本","~$1.3")],
    "現行 prod 架構 [best]",WHITE,GREEN,GREEN,
    "同一顆 GLM-5.2,跑在 Planner/Builder/Repair 重構後 builder(2026-07-02)。17/17 在 GLM 12-15 變異帶上緣;嚴謹解讀=零回歸,非架構提升品質。")

# ---- 8 verdict
s=slide(); header(s,"結論","定位與建議",BLUE)
items=[
 ("GLM-5.2 — prod 預設(不動)",GREEN,GREENBG,"嚴格品質與 Opus 僅差 1 案(xbar_r 也做對),成本最低、cache 唯一有效、速度最快 tier。CP 值完勝。"),
 ("Opus 4.8 — 難題升級層(on-demand)",PURP,PURPBG,"嚴格 16/17 最高。direct+cache 實測 ~20× GLM 成本。只在 fleet/composite/特定管制圖人工升級時叫用,且走 direct Anthropic 才有 cache。"),
 ("Sonnet 5 — 無角色",AMBER,AMBERBG,"嚴格 15/17(xbar_r 退化),direct+cache 成本仍 ~5× GLM。比 GLM 貴又不更準。"),
 ("KIMI K2.5 — 已淘汰",RED,REDBG,"最慢、cache 0%、品質不優於 GLM。"),
]
y=1.5
for title,acc,bg,body in items:
    rect(s,0.55,y,12.23,1.12,bg,acc,1.0)
    settext(box(s,0.85,y+0.12,11.6,0.9).text_frame,[
        para([seg(title,13.5,True,acc)],sa=3),para([seg(body,12,False,INK)],ls=1.05)],anchor=MSO_ANCHOR.MIDDLE)
    y+=1.26
settext(box(s,0.55,y+0.02,12.23,0.5).text_frame,[para([
    seg("共同待辦:",12.5,True,INK2),seg("apc-drift 的 weco_rules 四模型皆弱 → 修 knowledge / block doc,換模型無效。",12.5,False,INK)])])
footer(s,8)

out="docs/LLM_BAKEOFF_4MODEL.pptx"; prs.save(out)
print("saved:",out,"slides:",len(prs.slides._sldIdLst))
