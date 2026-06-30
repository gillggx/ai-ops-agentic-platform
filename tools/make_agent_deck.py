# -*- coding: utf-8 -*-
"""
技術分享 deck：AIOps 智能 Agent 平台 — 架構與設計
涵蓋：Agent 介紹 / Agent 行為(graph-first) / Context Engineering /
Memory Management / Block·Skill Document 設計 / 【新版】Skill→Auto Patrol·Check
資料來源：直接讀 code（skills_v2 V66、graph_build v30、agent_orchestrator_v2、
agent_knowledge、pb_blocks/seed.py）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x0B, 0x1F, 0x4D)
BLUE   = RGBColor(0x1D, 0x4E, 0xD8)
BLUEBG = RGBColor(0xEF, 0xF6, 0xFF)
INK    = RGBColor(0x1F, 0x29, 0x33)
INK2   = RGBColor(0x10, 0x18, 0x28)
MUT    = RGBColor(0x66, 0x70, 0x85)
LINE   = RGBColor(0xE4, 0xE7, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x06, 0x76, 0x47)
GREENBG= RGBColor(0xEC, 0xFD, 0xF3)
RED    = RGBColor(0xB4, 0x23, 0x18)
REDBG  = RGBColor(0xFE, 0xF3, 0xF2)
AMBER  = RGBColor(0xB5, 0x47, 0x08)
AMBERBG= RGBColor(0xFF, 0xFA, 0xEB)
GREYBG = RGBColor(0xF8, 0xFA, 0xFC)
SKY    = RGBColor(0x9D, 0xC0, 0xFF)
PURP   = RGBColor(0x6D, 0x28, 0xD9)
PURPBG = RGBColor(0xF5, 0xF3, 0xFF)

FONT = "Microsoft JhengHei"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(): return prs.slides.add_slide(BLANK)
def box(s, l, t, w, h): return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def rect(s, l, t, w, h, fill, line=None, line_w=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def rrect(s, l, t, w, h, fill, line=None, line_w=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = 0.08
    except Exception: pass
    return sp


def seg(t, sz=14, b=False, c=INK, i=False, mono=False):
    return {"t": t, "sz": sz, "b": b, "c": c, "i": i, "mono": mono}


def para(segs, align=PP_ALIGN.LEFT, sa=4, sb=0, ls=1.0):
    return {"segs": segs, "align": align, "space_after": sa, "space_before": sb, "line_sp": ls}


def settext(tf, runs, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, p in enumerate(runs):
        pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pp.alignment = p.get("align", PP_ALIGN.LEFT)
        pp.space_after = Pt(p.get("space_after", 4))
        pp.space_before = Pt(p.get("space_before", 0))
        pp.line_spacing = p.get("line_sp", 1.0)
        for sgmt in p["segs"]:
            r = pp.add_run(); r.text = sgmt["t"]
            f = r.font
            f.name = MONO if sgmt.get("mono") else FONT
            f.size = Pt(sgmt.get("sz", 14)); f.bold = sgmt.get("b", False)
            f.italic = sgmt.get("i", False); f.color.rgb = sgmt.get("c", INK)


def header(s, kicker, title, sub=None, accent=BLUE):
    rect(s, 0, 0, 13.333, 1.16, NAVY)
    rect(s, 0, 1.16, 13.333, 0.06, accent)
    t = box(s, 0.55, 0.13, 12.2, 1.0)
    runs = [para([seg(kicker, 11, True, SKY)], sa=2),
            para([seg(title, 22, True, WHITE)], sa=0)]
    if sub: runs.append(para([seg(sub, 11.5, False, RGBColor(0xCE, 0xDC, 0xF5))], sb=2))
    settext(t.text_frame, runs)


def footer(s, n, tag=""):
    t = box(s, 0.55, 7.04, 11, 0.33)
    settext(t.text_frame, [para([seg("AIOps 智能 Agent 平台 — 架構與設計分享" + (("  ·  " + tag) if tag else ""), 9, False, MUT)])])
    p = box(s, 12.2, 7.04, 0.7, 0.33)
    settext(p.text_frame, [para([seg(str(n), 9, True, MUT)], align=PP_ALIGN.RIGHT)])


def table(s, l, t, col_w, rows, header_fill=NAVY, row_h=0.34, fsz=11, hsz=11):
    total = sum(col_w); cur = t
    for ri, row in enumerate(rows):
        if ri == 0: rect(s, l, cur, total, row_h, header_fill)
        elif ri % 2 == 0: rect(s, l, cur, total, row_h, GREYBG)
        x = l
        for ci, cell in enumerate(row):
            opts = {}; txt = cell
            if isinstance(cell, tuple): txt, opts = cell
            tb = box(s, x + 0.03, cur, col_w[ci] - 0.06, row_h)
            c = WHITE if ri == 0 else opts.get("c", INK)
            b = True if ri == 0 else opts.get("b", False)
            sz = hsz if ri == 0 else opts.get("sz", fsz)
            al = opts.get("align", PP_ALIGN.LEFT)
            mono = opts.get("mono", False)
            settext(tb.text_frame, [para([seg(txt, sz, b, c, mono=mono)], align=al, ls=0.95)], anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[ci]
        cur += row_h
    for ri in range(len(rows) + 1):
        rect(s, l, t + ri * row_h - 0.005, total, 0.012, LINE)
    return cur


def divider(s, part, title, sub, n):
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 1.0, 3.0, 1.3, 0.12, BLUE)
    t = box(s, 1.0, 3.25, 11.2, 2)
    settext(t.text_frame, [
        para([seg(part, 15, True, SKY)], sa=8),
        para([seg(title, 30, True, WHITE)], sa=6),
        para([seg(sub, 14, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0)])
    footer(s, n)


# ===================================================== 1 TITLE
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.7, 13.333, 0.05, BLUE)
t = box(s, 0.9, 1.55, 11.6, 2.9)
settext(t.text_frame, [
    para([seg("技術分享 · Engineering Deep-Dive", 14, True, SKY)], sa=10),
    para([seg("AIOps 智能 Agent 平台", 40, True, WHITE)], sa=2),
    para([seg("架構、行為與設計原則", 24, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0)])
t2 = box(s, 0.9, 4.9, 11.6, 1.7)
settext(t2.text_frame, [
    para([seg("涵蓋主題：", 13, True, WHITE)], sa=6),
    para([seg("Agent 架構與行為　·　Context Engineering　·　Memory / Knowledge Management", 13, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=3),
    para([seg("Block / Skill Document 設計　·　【最新版】Skill → Auto Patrol / Auto Check 機制", 13, False, RGBColor(0xCE, 0xDC, 0xF5))], sa=0)])

# ===================================================== 2 全景
s = slide()
header(s, "全景", "平台與 Agent 的位置", "四個服務、清楚邊界 — 所有 Agent 都住在 Python sidecar，Java 是唯一 DB 擁有者")
svc = [
    ("aiops-app : 8000", "Next.js 前端", "只做 UI 渲染 + /api proxy；不直接碰 sidecar / DB", BLUE),
    ("java-api : 8002", "Spring Boot（唯一 DB 擁有者）", "所有 PostgreSQL 讀寫、auth、business CRUD、registry、alarms", GREEN),
    ("python-sidecar : 8050", "LangGraph Agents + Pipeline Executor", "三個 Agent surface + 57 blocks 全在這裡 in-process 跑", PURP),
    ("ontology-sim : 8012", "資料來源", "純資料服務，不知道 Agent 存在；介面同 production ontology", MUT),
]
y = 1.55
for title, role, desc, accent in svc:
    rrect(s, 0.55, y, 12.23, 1.18, WHITE, LINE, 1.0)
    rect(s, 0.55, y, 0.13, 1.18, accent)
    tb = box(s, 0.85, y + 0.13, 11.7, 0.95)
    settext(tb.text_frame, [
        para([seg(title + "   ", 14, True, accent), seg(role, 13, True, INK2)], sa=3),
        para([seg(desc, 12.5, False, INK)], sa=0, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.32
footer(s, 2, "全景")

# ===================================================== 3 三個 surface
s = slide()
header(s, "AGENT 介紹", "三個 Agent Surface", "同一個 sidecar、不同入口與編排器；都用 LangGraph，路由由 graph node 決定")
rows = [
    ["Surface", "入口 Endpoint", "編排器 Orchestrator", "路由方式"],
    ["Chat（維運對話）", ("/internal/agent/chat", {"mono": True, "sz": 10}), "agent_orchestrator_v2", "intent_classifier(5 bucket) + completeness gate"],
    ["Builder Glass Box（建 pipeline）", ("/internal/agent/build", {"mono": True, "sz": 10}), "graph_build v30", "goal_plan → phase_loop → verifier"],
    ["Block Advisor（block Q&A）", ("/internal/agent/build", {"mono": True, "sz": 10}), "advisor graph", "classify_advisor_intent 5 bucket dispatch"],
]
table(s, 0.55, 1.6, [3.5, 2.95, 2.7, 3.08], rows, row_h=0.62, fsz=11, hsz=11)
# bottom note
cb = rrect(s, 0.55, 4.35, 12.23, 2.15, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
settext(cb.text_frame, [
    para([seg("共同特徵", 14, True, INK2)], sa=6),
    para([seg("•  LLM 只在 node 內做窄任務", 13, True, BLUE), seg("（分類 / 抽參 / 規劃 / 寫答）；", 13, False, INK),
          seg("『下一步做什麼』永遠由 graph 的 conditional edge 決定。", 13, False, INK)], sa=5, ls=1.1),
    para([seg("•  Builder 入口先跑 ", 13, False, INK), seg("classify_advisor_intent", 12, True, PURP, mono=True),
          seg("：BUILD → Glass Box 建構；EXPLAIN / COMPARE / RECOMMEND → advisor 回答。", 13, False, INK)], sa=5, ls=1.1),
    para([seg("•  Prod builder LLM = ", 13, False, INK), seg("GLM-5.2 @ Fireworks @ medium", 13, True, GREEN),
          seg("（2026-06-25 經 SLASH-17 bake-off 由 KIMI 切換）。", 13, False, INK)], sa=0, ls=1.1),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3, "Agent 介紹")

# ===================================================== 4 核心哲學
s = slide()
header(s, "AGENT 行為 · 核心哲學", "Flow 在 Graph，LLM 只負責 Reason", accent=PURP)
# left principle
rrect(s, 0.55, 1.5, 5.55, 5.05, PURPBG, RGBColor(0xC4, 0xB5, 0xFD), 1.0)
tb = box(s, 0.85, 1.72, 5.0, 4.7)
settext(tb.text_frame, [
    para([seg("Hard Rule", 14, True, PURP)], sa=6),
    para([seg("任何『下一步 / 該呼叫哪個 tool』的決策，禁止塞進 LLM system prompt — 改寫成 graph node + deterministic dispatch。", 13, False, INK)], sa=12, ls=1.15),
    para([seg("為什麼", 14, True, PURP)], sa=6),
    para([seg("•  LLM 自由意志會違抗 prompt 規則（已多次證實）", 13, False, INK)], sa=7, ls=1.1),
    para([seg("•  prompt 寫的 flow 不可單測；graph node 是 pure function，可單測", 13, False, INK)], sa=7, ls=1.1),
    para([seg("•  出錯時 graph 知道是哪個 node fail；prompt-flow 只能猜『LLM 又走偏了』", 13, False, INK)], sa=0, ls=1.1),
])
# right three examples
rx = 6.35
examples = [
    ("意圖路由", "LLM 只輸出 enum bucket（5–7 種）；graph 的 conditional_edge 讀那個字串路由。intent 壞掉就 graceful fallback。"),
    ("Phase 驗證", "phase_verifier 是 deterministic 結構檢查（covers / orphan / leaf），不是問 LLM『這 phase 好了嗎』— 省錢且不會幻覺自評。"),
    ("修正預算", "repair / revise / reflect 都有 graph 控制的次數上限；何時放棄由 graph 決定，不靠 prompt 裡的信心分數。"),
]
y = 1.5
for title, body in examples:
    rrect(s, rx, y, 6.43, 1.55, WHITE, LINE, 1.0)
    rect(s, rx, y, 0.11, 1.55, PURP)
    tb = box(s, rx + 0.28, y + 0.14, 6.0, 1.3)
    settext(tb.text_frame, [
        para([seg(title, 13.5, True, PURP)], sa=4),
        para([seg(body, 12, False, INK)], sa=0, ls=1.08)], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.7
footer(s, 4, "Agent 行為")

# ===================================================== 5 Builder build graph
s = slide()
header(s, "AGENT 行為 · Builder", "v30 建構流程：規劃 → 逐 phase 建 → 驗證", "目標導向 ReAct：規劃只談意圖、不選 block；選 block 留到 execute 層")
# flow chain
nodes = [
    ("goal_plan", "出 3–7 個 phase\n只談意圖，禁列 block 名", BLUE),
    ("confirm gate", "interrupt\nuser 確認 / 編輯 phase", AMBER),
    ("phase_loop", "ReAct，1 action/round\n≤32 round/phase", PURP),
    ("phase_verifier", "結構檢查\ncovers / orphan / leaf", GREEN),
    ("finalize", "dry-run 執行\n輸出結果", NAVY),
]
x = 0.55; w = 2.3
for i, (title, body, accent) in enumerate(nodes):
    rrect(s, x, 1.65, w, 1.5, WHITE, accent, 1.5)
    rect(s, x, 1.65, w, 0.42, accent)
    hb = box(s, x + 0.1, 1.68, w - 0.2, 0.38)
    settext(hb.text_frame, [para([seg(title, 12.5, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    bb = box(s, x + 0.1, 2.12, w - 0.2, 0.98)
    settext(bb.text_frame, [para([seg(l, 10.5, False, INK)], align=PP_ALIGN.CENTER, ls=1.0) for l in body.split("\n")], anchor=MSO_ANCHOR.MIDDLE)
    if i < len(nodes) - 1:
        ar = box(s, x + w - 0.02, 2.05, 0.3, 0.7)
        settext(ar.text_frame, [para([seg("→", 18, True, MUT)], align=PP_ALIGN.CENTER)])
    x += w + 0.075
# phase kinds
pb = rrect(s, 0.55, 3.5, 12.23, 1.05, GREYBG, LINE, 1.0)
settext(pb.text_frame, [
    para([seg("Phase 的 7 種 expected output：", 13, True, INK2),
          seg("raw_data → transform → verdict → chart → table → scalar → alarm", 13, True, BLUE, mono=True)], sa=4, ls=1.1),
    para([seg("每個 phase 由一個 block『覆蓋』完成；phase_verifier 用 block 的 ", 12.5, False, INK),
          seg("produces.covers", 12, True, GREEN, mono=True), seg(" 配對 phase.expected。", 12.5, False, INK)], sa=0, ls=1.1),
], anchor=MSO_ANCHOR.MIDDLE)
# loop tools
lb = rrect(s, 0.55, 4.75, 12.23, 1.7, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
settext(lb.text_frame, [
    para([seg("phase_loop 裡 LLM 每 round 只發一個 tool call：", 13, True, INK2)], sa=5),
    para([seg("inspect_node_output", 12, True, BLUE, mono=True), seg(" 看上游資料　", 12, False, INK),
          seg("inspect_block_doc", 12, True, BLUE, mono=True), seg(" 讀 block 文件　", 12, False, INK),
          seg("add_node / connect / set_param", 12, True, BLUE, mono=True), seg(" 改 canvas　", 12, False, INK),
          seg("run_verifier", 12, True, BLUE, mono=True), seg(" 觸發驗證", 12, False, INK)], sa=5, ls=1.15),
    para([seg("param key 必須 100% 來自 param_schema（不准同義詞）；8 round 後（chart/alarm 12）verifier 自動觸發。", 12, False, INK)], sa=0, ls=1.1),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5, "Agent 行為")

# ===================================================== 6 self-correction + budgets
s = slide()
header(s, "AGENT 行為 · 自我修正", "卡住怎麼辦：有界修正 + 人工接手", "每條修正路徑都有 graph 控制的預算上限 — 防止 LLM 無限繞圈")
# left: escalation chain
rrect(s, 0.55, 1.55, 6.0, 4.95, WHITE, LINE, 1.0)
tb = box(s, 0.85, 1.75, 5.5, 4.6)
settext(tb.text_frame, [
    para([seg("升級路徑", 14, True, INK2)], sa=8),
    para([seg("1  phase_loop 跑完 round 仍沒過", 13, True, PURP)], sa=3),
    para([seg("     ↓", 12, False, MUT)], sa=3),
    para([seg("2  phase_revise（LLM 自省）", 13, True, AMBER)], sa=2),
    para([seg("     找根因 + 提 1 個替代策略 + 列能力缺口；只給 1 次", 11.5, False, INK)], sa=3, ls=1.05),
    para([seg("     ↓ 還是不行", 12, False, MUT)], sa=3),
    para([seg("3  halt_handover（interrupt 交給人）", 13, True, RED)], sa=2),
    para([seg("     user 選：edit_goal / take_over / backlog / abort", 11.5, False, INK)], sa=8, ls=1.05),
    para([seg("inspect_execution", 12, True, GREEN, mono=True), seg(" 另抓語意反模式：", 12.5, False, INK)], sa=3),
    para([seg("時序圖只有 <3 個 x 點 → 幾乎一定是上游被加了 limit=1，自動標記讓 reflect 修。", 11.5, False, INK)], sa=0, ls=1.1),
])
# right: budget table
rows = [
    ["預算 / 上限", "值"],
    ["graph recursion limit", ("60", {"mono": True})],
    ["ReAct round / phase", ("32", {"mono": True})],
    ["phase revise / phase", ("1", {"mono": True})],
    ["reflect_plan 次數", ("2", {"mono": True})],
    ["leaf prune 觸發", ("3 次 reject", {"mono": True})],
    ["fast-forward chain", ("4 phase", {"mono": True})],
    ["inspect / round", ("5", {"mono": True})],
]
table(s, 6.78, 1.55, [3.9, 2.1], rows, row_h=0.6, fsz=12, hsz=12)
footer(s, 6, "Agent 行為")

# ===================================================== 7 Context Eng 1
s = slide()
header(s, "CONTEXT ENGINEERING (1/2)", "每一 round 的 Observation 怎麼組", accent=GREEN)
# left sections
rrect(s, 0.55, 1.5, 6.0, 5.05, WHITE, LINE, 1.0)
rect(s, 0.55, 1.5, 6.0, 0.45, GREEN)
hb = box(s, 0.8, 1.54, 5.5, 0.38)
settext(hb.text_frame, [para([seg("13 段式 Observation（round 1 全量）", 13, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
tb = box(s, 0.85, 2.1, 5.5, 4.3)
secs = "已完成 phase｜當前 phase 目標｜verifier 回饋｜全 phase 脈絡｜可用 inputs｜canvas 節點 schema｜可接線選項｜近 6 步行動史｜user 指令｜matching blocks｜完整 catalog｜subphase 提示｜canvas 快照"
settext(tb.text_frame, [
    para([seg(secs, 12, False, INK)], sa=10, ls=1.25),
    para([seg("Round 2+ 只送 delta", 13, True, GREEN), seg("（canvas_diff + verifier 回饋 + phase 目標），訊息堆疊上限 32、指令截 600 字 — 把 token 壓住。", 12.5, False, INK)], sa=0, ls=1.15),
])
# right tiered catalog
rrect(s, 6.78, 1.5, 6.0, 5.05, GREENBG, RGBColor(0xA6, 0xF4, 0xC5), 1.0)
tb = box(s, 7.05, 1.72, 5.5, 4.7)
settext(tb.text_frame, [
    para([seg("分層 Catalog（目錄不是教科書）", 14, True, GREEN)], sa=8),
    para([seg("Tier 1 — 約 10 個核心 block", 13, True, INK2)], sa=3),
    para([seg("完整描述（≤1500 字）+ param_schema 直接 inline；agent 不用先 inspect 就能選。", 12, False, INK)], sa=10, ls=1.1),
    para([seg("Tier 2 — 其餘所有 block", 13, True, INK2)], sa=3),
    para([seg("只給一行索引；加進 canvas 前必須先 inspect_block_doc，否則 verifier 以『param 猜錯』退回。", 12, False, INK)], sa=10, ls=1.1),
    para([seg("matching_blocks", 12, True, GREEN, mono=True), seg(" 再依當前 phase.expected 的 covers 過濾，只給相關 block — 雙層收斂。", 12, False, INK)], sa=0, ls=1.1),
])
footer(s, 7, "Context Engineering")

# ===================================================== 8 Context Eng 2
s = slide()
header(s, "CONTEXT ENGINEERING (2/2)", "成本槓桿：Prompt Cache + 主動 RAG 查詢", accent=GREEN)
# left prompt cache
rrect(s, 0.55, 1.55, 6.0, 4.95, WHITE, LINE, 1.0)
tb = box(s, 0.85, 1.78, 5.5, 4.5)
settext(tb.text_frame, [
    para([seg("Prompt Cache", 15, True, GREEN)], sa=8),
    para([seg("那段大 observation prefix 每 round 重複 → cache 命中率 = 成本主槓桿。", 13, False, INK)], sa=10, ls=1.15),
    para([seg("•  system + tools + 最後一則訊息打上 ", 12.5, False, INK), seg("cache_control: ephemeral", 11.5, True, GREEN, mono=True)], sa=7, ls=1.1),
    para([seg("•  用 ", 12.5, False, INK), seg("OPENROUTER_PROVIDER_ORDER=Fireworks", 11.5, True, GREEN, mono=True), seg(" 釘住會吐 cache 的 provider", 12.5, False, INK)], sa=7, ls=1.1),
    para([seg("•  實測命中 ", 12.5, False, INK), seg("40–58%", 13, True, GREEN), seg("；後續 round 只付新增 token", 12.5, False, INK)], sa=7, ls=1.1),
    para([seg("•  provider 選錯差 ~8×（GLM DeepInfra 299s vs Fireworks 35s）", 12.5, False, INK)], sa=0, ls=1.1),
])
# right RAG
rrect(s, 6.78, 1.55, 6.0, 4.95, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
tb = box(s, 7.05, 1.78, 5.5, 4.5)
settext(tb.text_frame, [
    para([seg("從『推全清單』改『LLM 主動查』", 14, True, BLUE)], sa=8),
    para([seg("早期把整份 catalog 塞進 prompt → verbose catalog 反而誘發 over-build。", 13, False, INK)], sa=10, ls=1.15),
    para([seg("改為按需查詢工具：", 13, True, INK2)], sa=6),
    para([seg("query_blocks", 12, True, BLUE, mono=True), seg("　依語意找 block", 12.5, False, INK)], sa=6),
    para([seg("query_columns", 12, True, BLUE, mono=True), seg("　找欄位", 12.5, False, INK)], sa=6),
    para([seg("query_connectable_sources", 11.5, True, BLUE, mono=True), seg("　找可接來源", 12.5, False, INK)], sa=10),
    para([seg("效果：catalog 當『目錄』、深文件 lazy-load，prompt 精簡又夠用。", 12.5, False, INK)], sa=0, ls=1.1),
])
footer(s, 8, "Context Engineering")

# ===================================================== 9 Memory 1
s = slide()
header(s, "MEMORY / KNOWLEDGE (1/2)", "兩種知識來源 + 向量檢索", "知識存在 DB（Postgres + pgvector），不寫死在 prompt", accent=AMBER)
rows = [
    ["", "agent_knowledge", "block_docs"],
    ["粒度", "pipeline / plan 級規則", "單一 block 用法"],
    ["存放", "agent_knowledge 表 + 1024 維向量", "block_docs 表（markdown）"],
    ["欄位 / 結構", "applies_to · always_on · priority · scope", "frontmatter desc + body + column_docs"],
    ["檢索方式", "(1) always_on 高優先直推 (2) RAG cosine top-3", "inspect_block_doc 即時讀 DB"],
    ["用途", "形塑 phase 拆解、block 選擇規則", "catalog 一行摘要 + 完整參數文件"],
]
table(s, 0.55, 1.65, [2.0, 5.3, 4.93], rows, row_h=0.585, fsz=11, hsz=11.5)
cb = rrect(s, 0.55, 5.35, 12.23, 1.1, AMBERBG, RGBColor(0xFD, 0xE6, 0x8A), 1.0)
settext(cb.text_frame, [
    para([seg("Embedding：", 13, True, AMBER), seg("Cohere embed-multilingual-v3.0（1024 維）；pgvector 寫入用 native ", 12.5, False, INK),
          seg("CAST(:vec AS vector)", 11.5, True, AMBER, mono=True), seg("（JPA 不能隱式轉），背景每 30s 補沒 embedding 的列。", 12.5, False, INK)], ls=1.15),
    para([seg("為何要 always_on 直推：", 13, True, AMBER), seg("中文長查詢的向量召回偶爾會 miss，第一性原則（SPC=站級、APC=recipe 級）必須無條件到達。", 12.5, False, INK)], sb=4, ls=1.15),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 9, "Memory Management")

# ===================================================== 10 Memory 2
s = slide()
header(s, "MEMORY / KNOWLEDGE (2/2)", "分層注入：對的知識到對的 Agent 層", "Plan 層與 Execute 層各取各的（V58 / V59）", accent=AMBER)
# two layer cards
cards = [
    ("Plan 層（goal_plan，整段開頭一次）", BLUE, [
        "注入 always_on 核心 + plan-RAG top-3",
        "applies_to ∈ {plan, both}",
        "形塑意圖拆解（如『全廠查詢 → list_objects + foreach』）"]),
    ("Execute 層（phase_loop，每 phase 進入時）", PURP, [
        "只注入 execute-RAG top-3，不倒 always_on",
        "applies_to ∈ {execute, both}",
        "引導當下 phase 的 block / param 選擇"]),
]
x = 0.55; w = 6.0
for title, accent, items in cards:
    rrect(s, x, 1.55, w, 2.75, WHITE, LINE, 1.0)
    rect(s, x, 1.55, w, 0.5, accent)
    hb = box(s, x + 0.22, 1.59, w - 0.4, 0.42)
    settext(hb.text_frame, [para([seg(title, 12.5, True, WHITE)])], anchor=MSO_ANCHOR.MIDDLE)
    ib = box(s, x + 0.28, 2.2, w - 0.5, 2.0)
    settext(ib.text_frame, [para([seg("•  ", 13, True, accent), seg(it, 12.5, False, INK)], sa=9, ls=1.1) for it in items])
    x += w + 0.23
# why + session
wb = rrect(s, 0.55, 4.5, 6.0, 2.0, GREYBG, LINE, 1.0)
settext(wb.text_frame, [para([
    seg("為何要分層", 13, True, INK2)], sa=5),
    para([seg("goal_plan 設計上 block-agnostic（只出意圖）；『該選哪個 block』的知識若只在 plan 層讀，意圖被具體化後就遺失 — 必須送到真正選 block 的 phase_loop。", 12, False, INK)], sa=0, ls=1.15)],
    anchor=MSO_ANCHOR.MIDDLE)
sb = rrect(s, 6.78, 4.5, 6.0, 2.0, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
settext(sb.text_frame, [para([
    seg("對話記憶（Session）", 13, True, INK2)], sa=5),
    para([seg("agent_sessions", 12, True, BLUE, mono=True), seg(" 存跨 turn 的訊息歷史 + 累計 token + 最後 pipeline；知識則每 turn 從 DB 重新取（不快取在 session），確保拿到最新規則。", 12, False, INK)], sa=0, ls=1.15)],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10, "Memory Management")

# ===================================================== 11 Document SSoT
s = slide()
header(s, "BLOCK / SKILL DOCUMENT 設計", "Description = 唯一文件來源（1 寫、3 讀）", "改一個 DB 欄位，三個讀者立即同步，不需改 prompt、不需重部署")
# center write
rrect(s, 4.9, 1.6, 3.5, 0.95, NAVY)
tb = box(s, 4.95, 1.66, 3.4, 0.85)
settext(tb.text_frame, [
    para([seg("pb_blocks row（單一寫入）", 12.5, True, WHITE)], align=PP_ALIGN.CENTER, sa=2),
    para([seg("description · param_schema · examples", 10.5, False, SKY, mono=True)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
# three readers
readers = [
    ("Glass Box Agent", "建 pipeline 時讀 description 判斷『這 block 回傳什麼』、param_schema 驗參", BLUE, 0.55),
    ("Block Advisor", "回答『這 block 怎麼用 / A vs B』；參數段由 schema 直接 render，不讓 LLM 幻覺", PURP, 4.74),
    ("BlockDocsDrawer", "給人看的 UI；同一份 description / schema / examples 直接呈現", GREEN, 8.93),
]
for title, body, accent, x in readers:
    # connector
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 1.7), Inches(2.55), Inches(0.03), Inches(0.65))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    rrect(s, x, 3.2, 3.85, 2.0, WHITE, accent, 1.5)
    rect(s, x, 3.2, 3.85, 0.45, accent)
    hb = box(s, x + 0.15, 3.24, 3.55, 0.38)
    settext(hb.text_frame, [para([seg(title, 13, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    bb = box(s, x + 0.2, 3.75, 3.45, 1.35)
    settext(bb.text_frame, [para([seg(body, 12, False, INK)], ls=1.15)])
cb = rrect(s, 0.55, 5.5, 12.23, 0.95, AMBERBG, RGBColor(0xFD, 0xE6, 0x8A), 1.0)
settext(cb.text_frame, [para([
    seg("最高原則（CLAUDE.md #0）：", 13, True, AMBER),
    seg("LLM prompt 禁止 hardcode MCP / block 的用法、參數範例、回傳格式 — LLM 看不到 source code，只看得到 description。三邊不一致 = description 過時。", 12.5, False, INK)], ls=1.15)],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11, "Document 設計")

# ===================================================== 12 Block anatomy + invariant
s = slide()
header(s, "BLOCK / SKILL DOCUMENT 設計", "Block 結構、covers 配對、與 5 點不變式")
# left anatomy
rrect(s, 0.55, 1.55, 6.0, 4.95, WHITE, LINE, 1.0)
tb = box(s, 0.85, 1.75, 5.5, 4.6)
settext(tb.text_frame, [
    para([seg("pb_blocks 欄位", 14, True, INK2)], sa=7),
    para([seg("name · category · description · param_schema", 12, True, BLUE, mono=True)], sa=4, ls=1.1),
    para([seg("input/output_schema · examples · implementation", 12, True, BLUE, mono=True)], sa=10, ls=1.1),
    para([seg("produces.covers", 13, True, GREEN, mono=True), seg(" — 宣告此 block 產出哪種 phase 意圖", 12.5, False, INK)], sa=4),
    para([seg("{ raw_data · transform · chart · verdict · scalar · alarm }", 11.5, False, INK, mono=True)], sa=6, ls=1.05),
    para([seg("phase_verifier 用 covers 把 phase.expected 配對到正確 block；一個 block 可同時 cover 多種（如 threshold = verdict + scalar）。", 12.5, False, INK)], sa=0, ls=1.15),
])
# right invariant
rrect(s, 6.78, 1.55, 6.0, 4.95, GREYBG, LINE, 1.0)
tb = box(s, 7.05, 1.75, 5.5, 4.6)
settext(tb.text_frame, [
    para([seg("新增 block = 改 5 個地方", 14, True, INK2)], sa=7),
    para([seg("1  blocks/<name>.py　executor 實作", 12.5, False, INK)], sa=6),
    para([seg("2  __init__.py　BUILTIN_EXECUTORS 註冊", 12.5, False, INK)], sa=6),
    para([seg("3  real_executor.py　SIDECAR_NATIVE 白名單", 12.5, False, INK)], sa=6),
    para([seg("4  seed.py　catalog 定義", 12.5, False, INK)], sa=6),
    para([seg("5  Flyway V**.sql　DB row", 12.5, False, INK)], sa=10),
    para([seg("開機不變式", 13, True, GREEN), seg(" 印出『N builtin = native = seed = DB』必須相等（今約 57），缺一個就清楚報哪裡漏。", 12.5, False, INK)], sa=6, ls=1.15),
    para([seg("V54：System MCP 勾選即用 Haiku 自動生 block + skill 草稿（mcp_proxy），永遠當草稿，user review 後才 commit。", 12, False, INK)], sa=0, ls=1.1),
])
footer(s, 12, "Document 設計")

# ===================================================== 13 NEW skills_v2 model
s = slide()
header(s, "【最新版】SKILL → AUTO PATROL / AUTO CHECK", "新模型 skills_v2：一條 skill = 一條 pipeline + 自動化欄位",
       "V66（2026-06-28）取代舊的多步 skill_documents；V68 已 drop 舊表 — 這是目前 canonical 路徑", accent=RED)
# role derivation cards
roles = [
    ("tool", "無自動化", "純手動執行的 skill；trigger / gate / outcome 皆 NULL", MUT),
    ("patrol  = Auto Patrol", "自動化 + has_alarm", "pipeline 含 verdict 判斷式 → 條件達標會發 alarm", RED),
    ("datacheck = Auto Check", "自動化 + 無 alarm", "排程跑分析、產資料，但不發 alarm（outcome = data only）", BLUE),
]
x = 0.55; w = 4.0
for name, cond, body, accent in roles:
    rrect(s, x, 1.75, w, 1.95, WHITE, accent, 1.5)
    rect(s, x, 1.75, w, 0.5, accent)
    hb = box(s, x + 0.12, 1.79, w - 0.24, 0.42)
    settext(hb.text_frame, [para([seg(name, 12.5, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    bb = box(s, x + 0.18, 2.35, w - 0.36, 1.3)
    settext(bb.text_frame, [
        para([seg(cond, 12, True, accent)], align=PP_ALIGN.CENTER, sa=5),
        para([seg(body, 11.5, False, INK)], ls=1.1)], anchor=MSO_ANCHOR.MIDDLE)
    x += w + 0.115
# key facts
fb = rrect(s, 0.55, 3.95, 12.23, 2.55, GREYBG, LINE, 1.0)
settext(fb.text_frame, [
    para([seg("關鍵設計", 13.5, True, INK2)], sa=6),
    para([seg("•  role 是衍生的：", 13, False, INK), seg("has_alarm", 12, True, RED, mono=True),
          seg(" 由『pipeline 是否含 verdict node（block_step_check, isVerdict）』決定 — 沒判斷式就不能升 patrol。", 13, False, INK)], sa=6, ls=1.15),
    para([seg("•  自動化是同一 row 上的欄位：", 13, False, INK), seg("trigger_config", 12, True, BLUE, mono=True),
          seg(" (JSON: kind=schedule|event、schedule_spec、target、tool/source/event) · ", 12.5, False, INK),
          seg("alarm_gate", 12, True, BLUE, mono=True), seg(" · ", 12.5, False, INK), seg("outcome", 12, True, BLUE, mono=True),
          seg("。一條 SQL 取完，免 JOIN。", 13, False, INK)], sa=6, ls=1.15),
    para([seg("•  排程器永遠讀 ", 13, False, INK), seg("schedule_spec", 12, True, GREEN, mono=True),
          seg("（結構化），絕不 parse 顯示用的中文字串『每 1 小時』— flow 由結構決定，不靠字串解析。", 13, False, INK)], sa=0, ls=1.15),
], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 13, "Skill → Auto Patrol / Check【新版】")

# ===================================================== 14 NEW end-to-end flow
s = slide()
header(s, "【最新版】端到端資料流", "從『被撰寫的 skill』到『自動發出的 alarm』", accent=RED)
steps = [
    ("撰寫 + saveAutomation", "skills_v2 UI 設 role / trigger / gate；normalize 出 schedule_spec，寫成 1 row", BLUE),
    ("排程器 tick（每 60s）", "isDueV2 讀 schedule_spec 判斷到期；event 型則由上游 alarm 觸發", PURP),
    ("dispatch + fan-out", "『所有機台』且 pipeline 帶 $tool_id → 每台一個 run（{tool_id} 注入）", AMBER),
    ("Runner 執行 pipeline", "POST sidecar /internal/pipeline/execute；extractVerdict 掃出 pass:bool", GREEN),
    ("發 alarm + 事件擴散", "role=patrol & has_alarm & verdict=TRUE → AlarmEntity；再扇出給下游訂閱 skill", RED),
]
y = 1.6
for i, (title, body, accent) in enumerate(steps):
    rrect(s, 0.55, y, 8.4, 0.86, WHITE, LINE, 1.0)
    rect(s, 0.55, y, 0.7, 0.86, accent)
    nb = box(s, 0.55, y, 0.7, 0.86)
    settext(nb.text_frame, [para([seg(str(i + 1), 20, True, WHITE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    tb = box(s, 1.4, y + 0.08, 7.45, 0.72)
    settext(tb.text_frame, [
        para([seg(title, 12.5, True, accent)], sa=2),
        para([seg(body, 11, False, INK)], ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        ar = box(s, 4.4, y + 0.84, 0.5, 0.18)
        settext(ar.text_frame, [para([seg("↓", 13, True, MUT)], align=PP_ALIGN.CENTER)])
    y += 1.02
# right callouts
rrect(s, 9.2, 1.6, 3.58, 2.45, BLUEBG, RGBColor(0xBF, 0xDB, 0xFE), 1.0)
tb = box(s, 9.42, 1.78, 3.2, 2.15)
settext(tb.text_frame, [
    para([seg("追蹤", 13, True, BLUE)], sa=5),
    para([seg("每次執行都寫 ", 12, False, INK), seg("skill_runs", 11.5, True, BLUE, mono=True), seg(" 一列", 12, False, INK)], sa=5, ls=1.1),
    para([seg("沒發 alarm 時記 ", 12, False, INK), seg("alarm_skipped_reason", 11, True, BLUE, mono=True)], sa=5, ls=1.1),
    para([seg("→ patrol-activity 漏斗可逐關追", 12, False, INK)], sa=0, ls=1.1),
])
rrect(s, 9.2, 4.2, 3.58, 2.3, AMBERBG, RGBColor(0xFD, 0xE6, 0x8A), 1.0)
tb = box(s, 9.42, 4.38, 3.2, 2.0)
settext(tb.text_frame, [
    para([seg("語意（重要）", 13, True, AMBER)], sa=5),
    para([seg("verdict pass = TRUE", 12, True, RED, mono=True)], sa=4),
    para([seg("= 條件達標 = 有問題 = 發 alarm。", 12, False, INK)], sa=5, ls=1.1),
    para([seg("『檢查有符合』才是異常，不是正常。", 12, False, INK)], sa=0, ls=1.1),
])
footer(s, 14, "Skill → Auto Patrol / Check【新版】")

# ===================================================== 15 takeaways
s = slide()
header(s, "總結", "五個帶得走的設計原則")
items = [
    ("Flow 在 graph，LLM 只 reason", "下一步由 graph node 決定（可單測 / 可 trace）；LLM 只做分類、抽參、規劃、寫答。"),
    ("Document = 單一事實來源", "block / skill 的 description·schema·examples 一處寫、三處讀（Agent / Advisor / UI），不 drift、不重部署。"),
    ("Context engineering 兼顧連貫與成本", "13 段 observation + 分層 catalog + round delta + prompt cache（命中 40–58%）；catalog 是目錄不是教科書。"),
    ("知識分層注入", "agent_knowledge 依 applies_to 分 plan / execute；always_on 直推第一性原則，其餘走 pgvector RAG。"),
    ("skills_v2 把文件化 skill 變自動巡檢", "同一條 pipeline，加幾個欄位（role / trigger / gate）就成 Auto Patrol 或 Auto Check；has_alarm 由 verdict node 衍生。"),
]
y = 1.55
for i, (title, body) in enumerate(items):
    rrect(s, 0.55, y, 12.23, 0.95, WHITE if i % 2 else GREYBG, LINE, 1.0)
    nb = box(s, 0.7, y + 0.06, 0.7, 0.83)
    settext(nb.text_frame, [para([seg(str(i + 1), 22, True, BLUE)], align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    tb = box(s, 1.5, y + 0.1, 11.1, 0.78)
    settext(tb.text_frame, [
        para([seg(title, 13.5, True, INK2)], sa=2),
        para([seg(body, 12, False, INK)], ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02
footer(s, 15)

out = "docs/AGENT_ARCHITECTURE_SHARING.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
